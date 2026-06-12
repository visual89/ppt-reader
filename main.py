from fastapi import FastAPI, UploadFile, File
from pptx import Presentation
import tempfile

app = FastAPI()

@app.get("/")
def root():
    return {"status": "ok"}

def extract_shape_text(shape):
    texts = []

    # 일반 텍스트박스/도형
    if hasattr(shape, "text") and shape.text:
        texts.append(shape.text.strip())

    # 표 추출
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                cells.append(cell_text)
            texts.append(" | ".join(cells))

    # 그룹 도형 안쪽까지 추출
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            texts.extend(extract_shape_text(sub_shape))

    return texts

@app.post("/extract-ppt")
async def extract_ppt(file: UploadFile = File(...)):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name

    prs = Presentation(tmp_path)

    result = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            texts.extend(extract_shape_text(shape))

        cleaned_texts = []
        for t in texts:
            if t and t not in cleaned_texts:
                cleaned_texts.append(t)

        result.append({
            "slide": i,
            "text": "\n".join(cleaned_texts)
        })

    return result
